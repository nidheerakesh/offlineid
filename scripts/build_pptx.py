"""
build_pptx.py, generate the OfflineID hackathon deck as a themed .pptx.

Design language matches the app's "industrial biometric terminal" UI: near-black
background, signal-green accent, monospace readouts. 16:9, 16 slides.

Run:  .venv/Scripts/python.exe scripts/build_pptx.py
Out:  submission/OfflineID_Hackathon7.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---- palette (from src/ui/theme.ts) -------------------------------------------------
def C(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))


BG        = C("0A0E0D")
SURFACE   = C("111714")
SURFACE2  = C("161D1A")
LINE      = C("22302B")
LINEB     = C("2F423B")
ACCENT    = C("00E676")
ACCENTDIM = C("0B3D2A")
WARN      = C("FFB300")
DANGER    = C("FF5252")
INFO      = C("4FC3F7")
TEXT      = C("EAF3EF")
TEXTDIM   = C("8DA39B")
TEXTFAINT = C("5C726A")
ONACCENT  = C("04130C")

TITLE_FONT = "Segoe UI"
MONO       = "Consolas"

SW = 13.333
SH = 7.5
NSLIDES = 16

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---- primitives ---------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def hline(s, x, y, w, color=LINE, thick=0.012):
    return rect(s, x, y, w, thick, fill=color)


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def R(txt, size, color, bold=False, font=TITLE_FONT):
    return (txt, size, color, bold, font)


def chrome(s, idx, kicker=None, title=None):
    # left accent spine
    rect(s, 0, 0, 0.12, SH, fill=ACCENT)
    # footer: hairline + brand (with accent tab) + page number
    hline(s, 0.9, 6.88, 11.5)
    rect(s, 0.9, 6.99, 0.08, 0.17, fill=ACCENT)
    text(s, 1.1, 6.96, 8, 0.4,
         [[R("OFFLINE·ID", 9, TEXTFAINT, True, MONO),
           R("   NHAI HACKATHON 7.0 · DATALAKE 3.0", 9, TEXTFAINT, False, MONO)]])
    text(s, SW - 1.9, 6.96, 1.1, 0.4,
         [[R(f"{idx:02d} / {NSLIDES}", 9, TEXTFAINT, False, MONO)]], align=PP_ALIGN.RIGHT)
    y = 0.6
    if kicker:
        rect(s, 0.9, y + 0.03, 0.07, 0.2, fill=ACCENT)   # terminal-prompt tab
        text(s, 1.12, y, 11.3, 0.3, [[R(kicker, 12, ACCENT, True, MONO)]])
        y += 0.44
    if title:
        text(s, 0.88, y, 11.6, 1.0, [[R(title, 30, TEXT, True, TITLE_FONT)]])
        rect(s, 0.92, y + 0.74, 0.66, 0.05, fill=ACCENT)
    hline(s, 0.9, 1.82, 11.5)   # header separator
    return 2.02  # content top


def card(s, x, y, w, h, border=LINE, fill=SURFACE):
    return rect(s, x, y, w, h, fill=fill, line=border, line_w=1.25,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def card_label(s, x, y, w, label, color):
    """A titled card header: small accent tab + mono label."""
    rect(s, x + 0.22, y + 0.24, 0.32, 0.07, fill=color)
    text(s, x + 0.22, y + 0.4, w - 0.44, 0.4, [[R(label, 14, color, True, MONO)]])


def chip(s, x, y, label, color=ACCENT, w=2.0):
    c = rect(s, x, y, w, 0.44, fill=None, line=color, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.name = MONO
    r.font.color.rgb = color
    return c


def table(s, x, y, w, headers, rows, col_w, row_h=0.42, head_h=0.46, fs=13):
    nrows = len(rows) + 1
    ncols = len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w),
                            Inches(head_h + row_h * len(rows))).table
    gt.first_row = False
    gt.horz_banding = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for ri in range(nrows):
        gt.rows[ri].height = Inches(head_h if ri == 0 else row_h)
        for ci in range(ncols):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = ACCENTDIM
            else:
                cell.fill.fore_color.rgb = SURFACE if ri % 2 else SURFACE2
            val = headers[ci] if ri == 0 else rows[ri - 1][ci]
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(fs)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.name = MONO if (ri > 0 and ci > 0) else TITLE_FONT
            r.font.color.rgb = ACCENT if ri == 0 else (TEXT if ci == 0 else TEXTDIM)
    return gt


def bullets(s, x, y, w, items, fs=15, gap=10, mark="▸", mark_color=ACCENT, color=TEXT):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, sub = it
            paras.append([R(mark + "  ", fs, mark_color, True, TITLE_FONT),
                          R(head, fs, color, True, TITLE_FONT),
                          R("  " + sub, fs, TEXTDIM, False, TITLE_FONT)])
        else:
            paras.append([R(mark + "  ", fs, mark_color, True, TITLE_FONT),
                          R(it, fs, color, False, TITLE_FONT)])
    text(s, x, y, w, 5, paras, space_after=gap, line_spacing=1.05)


def flow_box(s, x, y, w, h, title, sub, accent=ACCENT):
    c = card(s, x, y, w, h, border=accent, fill=SURFACE)
    rect(s, x, y + 0.12, 0.07, h - 0.24, fill=accent)   # left colour tab
    text(s, x + 0.26, y + 0.1, w - 0.46, h - 0.2,
         [[R(title, 14, TEXT, True, TITLE_FONT)],
          [R(sub, 10.5, TEXTDIM, False, MONO)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    return c


def arrow_down(s, x, y, color=ACCENT):
    return rect(s, x, y, 0.32, 0.26, fill=color, shape=MSO_SHAPE.DOWN_ARROW)


def corner_frame(s):
    """Decorative bottom-right corner ticks (mirrors the top-left accent edges)."""
    rect(s, SW - 1.1, SH - 0.12, 1.1, 0.12, fill=ACCENT)
    rect(s, SW - 0.12, SH - 1.1, 0.12, 1.1, fill=ACCENT)


# =====================================================================================
#  SLIDE 1 - TITLE
# =====================================================================================
s = slide()
rect(s, 0, 0, 0.12, SH, fill=ACCENT)
rect(s, 0, 0, SW, 0.12, fill=ACCENT)
corner_frame(s)
rect(s, 0.9, 1.45, 0.08, 0.34, fill=ACCENT)
text(s, 1.12, 1.48, 11.5, 0.4, [[R("BIOMETRIC FIELD TERMINAL", 14, ACCENT, True, MONO)]])
text(s, 0.85, 2.0, 12, 1.6, [[R("OFFLINE·ID", 76, TEXT, True, TITLE_FONT)]])
rect(s, 0.92, 3.45, 2.4, 0.06, fill=ACCENT)
text(s, 0.9, 3.7, 11.0, 1.2,
     [[R("Secure offline facial recognition & liveness detection", 22, TEXTDIM, False, TITLE_FONT)],
      [R("for field personnel in zero-network zones.", 22, TEXTDIM, False, TITLE_FONT)]],
     space_after=2)
chip(s, 0.9, 5.2, "100% ON-DEVICE", ACCENT, 2.4)
chip(s, 3.45, 5.2, "ANDROID + iOS", INFO, 2.1)
chip(s, 5.7, 5.2, "9.1 MB MODELS", WARN, 2.2)
chip(s, 8.05, 5.2, "OPEN SOURCE", TEXTDIM, 2.0)
hline(s, 0.92, 6.35, 3.2, color=LINEB)
text(s, 0.9, 6.5, 11, 0.5,
     [[R("NHAI Hackathon 7.0", 13, TEXT, True, TITLE_FONT),
       R("   ·   Datalake 3.0 integration module   ·   v1.0.0", 13, TEXTFAINT, False, TITLE_FONT)]])

# =====================================================================================
#  SLIDE 2 - PROBLEM
# =====================================================================================
s = slide()
top = chrome(s, 2, "THE PROBLEM", "Authenticating field staff where there is no network")
bullets(s, 0.9, top, 11.6, [
    ("Zero-network sites.", "NHAI personnel work in remote zones; cloud face APIs simply fail offline."),
    ("Fraud risk.", "Photos and phone-screen replays let people fake attendance."),
    ("Constrained hardware.", "Must run on mid-range phones: 3 GB RAM, Android 8+/iOS 12+, no GPU."),
    ("Hard targets.", "< 1 second, > 95% accuracy, diverse Indian faces, harsh outdoor light."),
    ("Must fit Datalake 3.0.", "Ship as a React Native module, not a separate app."),
], fs=16, gap=12)
card(s, 0.9, top + 3.55, 11.5, 0.95, border=LINEB, fill=SURFACE2)
rect(s, 0.9, top + 3.55, 0.08, 0.95, fill=ACCENT)
text(s, 1.2, top + 3.7, 11.0, 0.7,
     [[R("“", 20, ACCENT, True, TITLE_FONT),
       R("How do we authenticate field personnel securely, fully offline, on standard "
         "phones, and integrate it into a React Native app on Android and iOS?", 14.5, TEXT, False, TITLE_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================================
#  SLIDE 3 - SOLUTION
# =====================================================================================
s = slide()
top = chrome(s, 3, "OUR SOLUTION", "A fully on-device authentication module")
text(s, 0.9, top, 11.6, 1.1,
     [[R("Detect → prove liveness → recognise in ", 18, TEXT, False, TITLE_FONT),
       R("~51 ms", 18, ACCENT, True, MONO),
       R(", with no internet, then ", 18, TEXT, False, TITLE_FONT),
       R("sync-and-purge", 18, ACCENT, True, TITLE_FONT),
       R(" to AWS when back online.", 18, TEXT, False, TITLE_FONT)]],
     line_spacing=1.1)
cards = [
    ("OFFLINE", "Recognition + liveness run\nentirely on the device.", ACCENT),
    ("ANTI-SPOOF", "Passive FASNet + active\ngesture defeat photos.", INFO),
    ("ENCRYPTED", "AES-256-GCM faceprints;\nno raw images stored.", WARN),
    ("SYNC & PURGE", "Presigned S3 upload, then\nlocal records deleted.", ACCENT),
]
cw = 2.78
gapx = 0.13
x0 = 0.9
for i, (h, b, col) in enumerate(cards):
    x = x0 + i * (cw + gapx)
    card(s, x, top + 1.45, cw, 2.55, border=col, fill=SURFACE)
    rect(s, x + 0.22, top + 1.68, 0.3, 0.07, fill=col)
    text(s, x + 0.22, top + 1.84, cw - 0.4, 2.0,
         [[R(h, 15, col, True, MONO)],
          [R(b, 13, TEXTDIM, False, TITLE_FONT)]], space_after=8, line_spacing=1.05)

# =====================================================================================
#  SLIDE 4 - PIPELINE
# =====================================================================================
s = slide()
top = chrome(s, 4, "THE AI PIPELINE", "Camera still to attendance record")
steps = [
    ("1 · Detect", "SCRFD-500M → box + 5 landmarks", ACCENT),
    ("2 · Passive liveness", "FASNet ×2 (2.7 / 4.0) anti-spoof", INFO),
    ("3 · Active liveness", "ML Kit gesture: blink / smile / turn", INFO),
    ("4 · Recognise", "ArcFace align → MobileFaceNet 512-d", ACCENT),
    ("5 · Match & log", "cosine vs enrolled → encrypted row", ACCENT),
    ("6 · Sync & purge", "reconnect → S3 PUT → local delete", WARN),
]
bx = 0.9
bw = 11.5
bh = 0.56
y = top + 0.02
for i, (t, sub, col) in enumerate(steps):
    flow_box(s, bx, y, bw, bh, t, sub, accent=col)
    if i < len(steps) - 1:
        arrow_down(s, bx + bw / 2 - 0.16, y + bh + 0.02, color=col)
    y += bh + 0.245

# =====================================================================================
#  SLIDE 5 - WHY INNOVATIVE
# =====================================================================================
s = slide()
top = chrome(s, 5, "INNOVATION · 30 MARKS", "Edge AI that earns its footprint")
trip = [
    ("9.1 MB", "total model bundle", "INT8-quantised MobileFaceNet keeps the whole\nbundle under half the 20 MB budget."),
    ("2 layers", "of liveness", "Passive FASNet anti-spoof PLUS a randomised\nactive gesture sequence."),
    ("0 cloud", "calls for auth", "Detection, liveness and matching all run\non-device; works in airplane mode."),
]
cw = 3.75
x0 = 0.9
for i, (big, small, body) in enumerate(trip):
    x = x0 + i * (cw + 0.18)
    card(s, x, top + 0.15, cw, 3.55, border=LINEB, fill=SURFACE)
    rect(s, x + 0.25, top + 0.42, 0.5, 0.07, fill=ACCENT)
    text(s, x + 0.25, top + 0.58, cw - 0.5, 1.3,
         [[R(big, 38, ACCENT, True, TITLE_FONT)],
          [R(small.upper(), 12, TEXTDIM, True, MONO)]], space_after=4)
    text(s, x + 0.25, top + 2.05, cw - 0.5, 1.5,
         [[R(body, 13.5, TEXT, False, TITLE_FONT)]], line_spacing=1.1)

# =====================================================================================
#  SLIDE 6 - MODEL STACK
# =====================================================================================
s = slide()
top = chrome(s, 6, "MODEL STACK", "Four open-source ONNX models")
table(s, 0.9, top + 0.12, 11.5,
      ["Stage", "Model", "Size", "Role"],
      [["Detect", "SCRFD-500M", "2.41 MB", "face box + 5 landmarks"],
       ["Recognise", "MobileFaceNet INT8", "3.35 MB", "512-d ArcFace embedding"],
       ["Liveness A", "MiniFASNet V2 (2.7)", "1.66 MB", "passive anti-spoof"],
       ["Liveness B", "MiniFASNet V1SE (4.0)", "1.66 MB", "passive anti-spoof"],
       ["Runtime", "ONNX Runtime Mobile", "~3.5 MB", "CPU / XNNPACK / NNAPI / CoreML"]],
      col_w=[2.0, 3.4, 1.7, 4.4], row_h=0.6, head_h=0.5, fs=13.5)
text(s, 0.9, top + 4.05, 11.5, 0.6,
     [[R("All MIT-licensed. ", 14, ACCENT, True, TITLE_FONT),
       R("Models 9.1 MB + ONNX Runtime ≈ 12.6 MB added to the Datalake APK, "
         "inside the 20 MB brief budget.", 14, TEXTDIM, False, TITLE_FONT)]],
     line_spacing=1.1)

# =====================================================================================
#  SLIDE 7 - LIVENESS / ANTI-SPOOF
# =====================================================================================
s = slide()
top = chrome(s, 7, "LIVENESS · ANTI-SPOOF", "Two independent layers, both must pass")
card(s, 0.9, top + 0.12, 5.65, 3.55, border=INFO, fill=SURFACE)
card_label(s, 0.9, top + 0.12, 5.65, "PASSIVE", INFO)
bullets(s, 1.15, top + 1.05, 5.2, [
    "FASNet two-scale crops (2.7 + 4.0)",
    "Softmax P(real) > 0.6 to pass",
    "Blocks printed photos, screens, 2D masks",
    "Runs in a few milliseconds",
], fs=13.5, gap=9, mark_color=INFO)
card(s, 6.75, top + 0.12, 5.65, 3.55, border=ACCENT, fill=SURFACE)
card_label(s, 6.75, top + 0.12, 5.65, "ACTIVE", ACCENT)
bullets(s, 7.0, top + 1.05, 5.2, [
    "Randomised gesture in a 5 s window",
    "Blink (EAR) · turn (yaw ±20°) · smile",
    "On-device ML Kit, no network",
    "Ordered sequence defeats replays",
], fs=13.5, gap=9, mark_color=ACCENT)
card(s, 0.9, top + 3.95, 11.5, 0.62, border=LINEB, fill=SURFACE2)
text(s, 1.2, top + 3.95, 11.0, 0.62,
     [[R("GATE   ", 14, WARN, True, MONO),
       R("recognition only runs after both liveness layers pass.", 14, TEXT, False, TITLE_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================================
#  SLIDE 8 - SECURITY
# =====================================================================================
s = slide()
top = chrome(s, 8, "SECURITY", "Privacy-first by construction")
bullets(s, 0.9, top + 0.1, 11.6, [
    ("Encrypted at rest.", "Faceprints sealed with AES-256-GCM; key in the Android Keystore."),
    ("No raw images.", "Only 512-d embeddings are stored, never photos of faces."),
    ("No cloud credentials.", "Sync uses short-TTL presigned S3 URLs; device holds no AWS keys."),
    ("Abuse resistant.", "30-second lockout after repeated failed attempts."),
    ("Purge on sync.", "Attendance rows are deleted locally once the upload is confirmed."),
], fs=16, gap=14)

# =====================================================================================
#  SLIDE 9 - OFFLINE DATA + SYNC
# =====================================================================================
s = slide()
top = chrome(s, 9, "OFFLINE-FIRST DATA", "Works fully offline, reconciles when online")
card(s, 0.9, top + 0.12, 5.65, 3.45, border=LINEB, fill=SURFACE)
card_label(s, 0.9, top + 0.12, 5.65, "ON DEVICE", ACCENT)
bullets(s, 1.15, top + 1.0, 5.2, [
    "SQLite: face_embeddings (encrypted, permanent)",
    "SQLite: attendance_log (ephemeral queue)",
    "Enrol · authenticate · log, all offline",
], fs=13.5, gap=10)
card(s, 6.75, top + 0.12, 5.65, 3.45, border=LINEB, fill=SURFACE)
card_label(s, 6.75, top + 0.12, 5.65, "SYNC & PURGE", WARN)
bullets(s, 7.0, top + 1.0, 5.2, [
    "NetInfo detects reconnect",
    "Batch ≤ 10 → presigned S3 PUT",
    "On 200 → delete local row (purge)",
    "Exponential backoff; 403 → refresh URL",
], fs=13.5, gap=10, mark_color=WARN)
text(s, 0.9, top + 3.78, 11.5, 0.4,
     [[R("Devices stay lean; the server becomes the single source of truth.", 13.5, TEXTDIM, False, TITLE_FONT)]])

# =====================================================================================
#  SLIDE 10 - PERFORMANCE
# =====================================================================================
s = slide()
top = chrome(s, 10, "PERFORMANCE · BENCHMARKS", "Measured against every target")
metrics = [("51 ms", "host-CPU pipeline", "target < 1 s"),
           ("9.1 MB", "model bundle", "cap 20 MB"),
           ("99.5%", "LFW accuracy", "target > 95%"),
           ("100%", "offline capable", "zero network")]
cw = 2.78
x0 = 0.9
for i, (big, small, tgt) in enumerate(metrics):
    x = x0 + i * (cw + 0.13)
    card(s, x, top + 0.1, cw, 1.95, border=LINEB, fill=SURFACE)
    text(s, x + 0.22, top + 0.28, cw - 0.4, 1.6,
         [[R(big, 30, ACCENT, True, TITLE_FONT)],
          [R(small.upper(), 11, TEXT, True, MONO)],
          [R(tgt, 11, TEXTFAINT, False, MONO)]], space_after=3)
table(s, 0.9, top + 2.3, 11.5,
      ["Model", "Size", "Avg", "P95"],
      [["SCRFD-500M (detect)", "2.41 MB", "6.1 ms", "7.5 ms"],
       ["MobileFaceNet INT8 (embed)", "3.35 MB", "43.6 ms", "49.5 ms"],
       ["MiniFASNet 2.7 / 4.0 (liveness)", "1.66 MB", "1.5 ms", "2.1 ms"]],
      col_w=[4.6, 2.0, 2.55, 2.35], row_h=0.44, head_h=0.46, fs=12.5)
text(s, 0.9, top + 4.2, 11.6, 0.8,
     [[R("Host CPU · ONNX Runtime.  Pipeline = detect + 2× liveness + embed; "
         "add align & cosine < 10 ms for end-to-end.", 12, TEXTDIM, False, TITLE_FONT)],
      [R("App footprint on Datalake ≈ 12.6 MB (9.1 MB models + 3.5 MB runtime).  "
         "INT8 vs FP32 cosine > 0.99 on aligned faces.", 12, TEXTDIM, False, TITLE_FONT)]],
     space_after=3, line_spacing=1.05)

# =====================================================================================
#  SLIDE 11 - ARCHITECTURE
# =====================================================================================
s = slide()
top = chrome(s, 11, "ARCHITECTURE", "React Native UI over a native ONNX engine")
layers = [
    ("UI · React Native + TypeScript", "Scan · Enrol · People · Sync · System screens", ACCENT),
    ("Orchestration · useFaceAuth hook", "state machine: detect → liveness → recognise", INFO),
    ("Services (JS)", "FaceEngine bridge · Liveness · Stores · Sync", INFO),
    ("Native engine · Kotlin (Android) / Swift (iOS)", "ONNX Runtime: detect · liveness · embed", ACCENT),
    ("Storage · SQLite + Keystore + AWS S3", "encrypted faceprints · attendance queue", WARN),
]
y = top + 0.05
bw = 11.5
bh = 0.74
for i, (t, sub, col) in enumerate(layers):
    flow_box(s, 0.9, y, bw, bh, t, sub, accent=col)
    if i < len(layers) - 1:
        arrow_down(s, 0.9 + bw / 2 - 0.16, y + bh - 0.02, color=col)
    y += bh + 0.22

# =====================================================================================
#  SLIDE 12 - CROSS-PLATFORM & INTEGRATION
# =====================================================================================
s = slide()
top = chrome(s, 12, "FEASIBILITY · 30 MARKS", "Drops into Datalake 3.0")
bullets(s, 0.9, top + 0.05, 11.6, [
    ("Register one native package.", "add FaceEnginePackage() (Android) / link the Swift module (iOS)."),
    ("Add the ONNX dependency + 4 model assets.", "≈ 12.6 MB total."),
    ("Import the screens & services.", "Auth / Enrol / Sync mount in the existing navigator."),
    ("Point SYNC_BASE_URL at a presigned endpoint.", "the only backend touchpoint."),
    ("Zero changes to the offline auth path.", "no other Datalake backend work."),
], fs=15.5, gap=11)
card(s, 0.9, top + 3.5, 11.5, 1.0, border=ACCENT, fill=SURFACE2)
rect(s, 0.9, top + 3.5, 0.08, 1.0, fill=ACCENT)
text(s, 1.2, top + 3.62, 11.0, 0.8,
     [[R("iOS engine is written ", 14, TEXT, False, TITLE_FONT),
       R("(ios/FaceEngine/, a 1:1 Swift port, same models & math)", 14, ACCENT, True, MONO),
       R("; remaining work is one-time Xcode build wiring on a Mac.", 14, TEXT, False, TITLE_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# =====================================================================================
#  SLIDE 13 - TECH & QUALITY
# =====================================================================================
s = slide()
top = chrome(s, 13, "ENGINEERING QUALITY", "Open-source, typed, tested")
bullets(s, 0.9, top + 0.1, 11.6, [
    ("Open-source only.", "MIT / Apache stack, no paid licences, source shared."),
    ("TypeScript strict.", "typecheck clean; 15 unit tests passing."),
    ("Standalone offline APK.", "release build embeds JS; runs in airplane mode, no Metro."),
    ("Reproducible models.", "Python export + validation scripts in the repo."),
    ("Documented.", "architecture, model pipeline, benchmarks, Datalake integration guide."),
], fs=16, gap=13)

# =====================================================================================
#  SLIDE 14 - ROADMAP
# =====================================================================================
s = slide()
top = chrome(s, 14, "SCALABILITY · SUSTAINABILITY · 20 MARKS", "Built to grow")
bullets(s, 0.9, top + 0.1, 11.6, [
    ("Scales to N users.", "cosine match < 5 ms at 500 enrolled faces."),
    ("Lean by design.", "SQLite queue → batch S3 → purge keeps storage flat."),
    ("Swappable models.", "update accuracy by replacing ONNX files, 8 MB headroom."),
    ("Demographic tuning.", "on-device fine-tune on an Indian face subset (roadmap)."),
    ("iOS device build + GPU delegate.", "next on the path to full parity."),
], fs=16, gap=13)

# =====================================================================================
#  SLIDE 15 - EVALUATION MAPPING
# =====================================================================================
s = slide()
top = chrome(s, 15, "SCORING", "How OfflineID maps to the rubric")
table(s, 0.9, top + 0.25, 11.5,
      ["Criterion", "Marks", "Where we win"],
      [["Innovation", "30", "9.1 MB edge AI · two-layer liveness"],
       ["Feasibility", "30", "drops into Datalake 3.0 · ~51 ms · < 1 s"],
       ["Scalability & Sustainability", "20", "offline queue · sync/purge · model swaps"],
       ["Presentation & Documentation", "20", "source clarity · integration guide · this deck"]],
      col_w=[4.6, 1.4, 5.5], row_h=0.68, head_h=0.52, fs=13.5)
text(s, 0.9, top + 3.95, 11.5, 0.5,
     [[R("Open-source · lightweight · fully offline · ready for Datalake 3.0.", 14, TEXTDIM, False, TITLE_FONT)]])

# =====================================================================================
#  SLIDE 16 - CLOSING
# =====================================================================================
s = slide()
rect(s, 0, 0, 0.12, SH, fill=ACCENT)
rect(s, 0, SH - 0.12, SW, 0.12, fill=ACCENT)
corner_frame(s)
rect(s, 0.9, 1.65, 0.08, 0.34, fill=ACCENT)
text(s, 1.12, 1.68, 11.5, 0.4, [[R("OFFLINE·ID", 16, ACCENT, True, MONO)]])
text(s, 0.85, 2.2, 11.6, 1.5,
     [[R("Authenticate anyone, anywhere,", 34, TEXT, True, TITLE_FONT)],
      [R("zero network required.", 34, TEXT, True, TITLE_FONT)]], space_after=2)
rect(s, 0.92, 3.7, 2.0, 0.05, fill=ACCENT)
text(s, 0.9, 4.0, 11.5, 1.2,
     [[R("Secure · lightweight · fully offline face auth that plugs into Datalake 3.0.",
         16, TEXTDIM, False, TITLE_FONT)]])
text(s, 0.9, 5.05, 11.5, 1.2,
     [[R("Source    ", 13, TEXTFAINT, True, MONO),
       R("github.com/moneytosms/offlineid", 13, ACCENT, False, MONO)],
      [R("Release   ", 13, TEXTFAINT, True, MONO),
       R("github.com/moneytosms/offlineid/releases/tag/v1.0.0", 13, ACCENT, False, MONO)]],
     space_after=6, line_spacing=1.2)
text(s, 0.9, 6.45, 11.5, 0.5, [[R("Thank you   ·   Q & A", 15, TEXT, True, TITLE_FONT)]])

# ---- save ---------------------------------------------------------------------------
import os
out = os.path.join("submission", "OfflineID_Hackathon7.pptx")
prs.save(out)
print("wrote", out, f"({len(prs.slides._sldIdLst)} slides)")
