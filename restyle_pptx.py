"""Restyle OfflineID_Hackathon7.pptx to Deep Navy / Sky Blue theme."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# Color palette
BG_DEEP       = RGBColor(0x08, 0x0C, 0x1A)  # #080C1A
BG_SURFACE    = RGBColor(0x0D, 0x12, 0x26)  # #0D1226
ACCENT_BLUE   = RGBColor(0x38, 0xBD, 0xF8)  # #38BDF8 sky blue
TEXT_PRIMARY  = RGBColor(0xE2, 0xE8, 0xF0)  # #E2E8F0
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)  # #94A3B8
NHAI_GREEN    = RGBColor(0x22, 0xC5, 0x5E)  # #22C55E keep brand
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)  # #F59E0B keep warn
NAVY_BORDER   = RGBColor(0x1E, 0x28, 0x47)  # #1E2847
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)

# Colors to keep as-is (accent, green, amber)
KEEP_COLORS = {
    (0x38, 0xBD, 0xF8),  # sky blue
    (0x22, 0xC5, 0x5E),  # NHAI green
    (0xF5, 0x9E, 0x0B),  # amber
    (0x16, 0xA3, 0x4A),  # darker green variant
    (0x15, 0x80, 0x3D),  # another green
    (0x4A, 0xDE, 0x80),  # light green
}

def is_light_color(r, g, b):
    """Return True if color is perceptually light (luminance > 0.4)."""
    lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return lum > 0.4

def is_green(r, g, b):
    """Return True if color is greenish."""
    return g > r and g > b and g > 80

def is_blue_accent(r, g, b):
    """Return True if color is a bright blue/cyan."""
    return b > r and b > g * 0.6 and (b > 150 or (r < 100 and g > 150))

def should_keep(r, g, b):
    """Return True if this color should be preserved."""
    return (r, g, b) in KEEP_COLORS

def map_fill_color(r, g, b):
    """Map a fill color to the new palette."""
    if should_keep(r, g, b):
        return RGBColor(r, g, b)
    if is_green(r, g, b):
        return NHAI_GREEN
    if is_light_color(r, g, b):
        return BG_SURFACE
    # Dark color → deep background
    return BG_DEEP

def map_text_color(r, g, b):
    """Map a text color to the new palette."""
    if should_keep(r, g, b):
        return RGBColor(r, g, b)
    if is_green(r, g, b):
        return NHAI_GREEN
    if is_blue_accent(r, g, b):
        return ACCENT_BLUE
    # White / very light → primary text
    if is_light_color(r, g, b) and r > 180 and g > 180 and b > 180:
        return TEXT_PRIMARY
    if is_light_color(r, g, b):
        return TEXT_PRIMARY
    # Dark / black text → primary text (invert for dark bg)
    return TEXT_PRIMARY

def map_line_color(r, g, b):
    """Map a line/border color to the new palette."""
    if should_keep(r, g, b):
        return RGBColor(r, g, b)
    if is_green(r, g, b):
        return ACCENT_BLUE  # convert green lines to sky blue
    if is_blue_accent(r, g, b):
        return ACCENT_BLUE
    if is_light_color(r, g, b):
        return NAVY_BORDER
    return NAVY_BORDER

def set_slide_background(slide, color):
    """Set slide background to solid fill color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    # Access the slide background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def restyle_run(run):
    """Restyle a text run's font color."""
    font = run.font
    try:
        if font.color and font.color.type is not None:
            try:
                rgb = font.color.rgb
                r, g, b = rgb[0], rgb[1], rgb[2]
                new_color = map_text_color(r, g, b)
                font.color.rgb = new_color
            except Exception:
                font.color.rgb = TEXT_PRIMARY
        else:
            # No explicit color set — set to primary text
            font.color.rgb = TEXT_PRIMARY
    except Exception:
        pass

def restyle_shape(shape):
    """Restyle a single shape."""
    # --- Fill ---
    try:
        fill = shape.fill
        from pptx.enum.dml import MSO_FILL
        if fill.type is not None:
            if fill.type.name == 'SOLID' or str(fill.type) == 'SOLID (1)':
                try:
                    rgb = fill.fore_color.rgb
                    r, g, b = rgb[0], rgb[1], rgb[2]
                    new_fill = map_fill_color(r, g, b)
                    fill.solid()
                    fill.fore_color.rgb = new_fill
                except Exception:
                    fill.solid()
                    fill.fore_color.rgb = BG_SURFACE
    except Exception:
        pass

    # --- Line ---
    try:
        line = shape.line
        if line.color and line.color.type is not None:
            try:
                rgb = line.color.rgb
                r, g, b = rgb[0], rgb[1], rgb[2]
                new_line = map_line_color(r, g, b)
                line.color.rgb = new_line
            except Exception:
                line.color.rgb = NAVY_BORDER
    except Exception:
        pass

    # --- Text frame ---
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    restyle_run(run)
    except Exception:
        pass

    # --- Table ---
    try:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    # Cell fill
                    try:
                        cf = cell.fill
                        if cf.type is not None:
                            try:
                                rgb = cf.fore_color.rgb
                                r, g, b = rgb[0], rgb[1], rgb[2]
                                new_fill = map_fill_color(r, g, b)
                                cf.solid()
                                cf.fore_color.rgb = new_fill
                            except Exception:
                                cf.solid()
                                cf.fore_color.rgb = BG_SURFACE
                    except Exception:
                        pass
                    # Cell text
                    try:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                restyle_run(run)
                    except Exception:
                        pass
    except Exception:
        pass

    # --- Group shapes (recurse) ---
    try:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for s in shape.shapes:
                restyle_shape(s)
    except Exception:
        pass


def main():
    pptx_path = r"D:\Projects\offlineid\submission\OfflineID_Hackathon7.pptx"
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx + 1} / {len(prs.slides)}")

        # Set slide background
        set_slide_background(slide, BG_SURFACE)

        # Restyle each shape
        for shape in slide.shapes:
            restyle_shape(shape)

    prs.save(pptx_path)
    print(f"Saved: {pptx_path}")


if __name__ == "__main__":
    main()
